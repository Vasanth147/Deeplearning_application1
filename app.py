import streamlit as st
from pathlib import Path
import os
import pdfplumber
from pptx import Presentation
import docx
import pandas as pd
from PIL import Image
import pytesseract
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
import uuid


st.set_page_config(page_title="Document Q&A RAG App", layout="wide")
st.title("Document Q&A RAG Application")
st.markdown("""
<style>
    .highlight { background-color: #ffff99; }
    .source-link { color: #1f77b4; text-decoration: none; font-weight: bold; }
    .source-link:hover { text-decoration: underline; }
</style>
""", unsafe_allow_html=True)


def extract_text_from_file(uploaded_file):
    """Extract text from supported file types."""
    suffix = uploaded_file.name.split(".")[-1].lower()
    text = ""

    if suffix == "pdf":
        with pdfplumber.open(uploaded_file) as pdf:
            for page in pdf.pages:
                text += page.extract_text() + "\n"
    elif suffix == "pptx":
        prs = Presentation(uploaded_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif suffix == "docx":
        doc = docx.Document(uploaded_file)
        for para in doc.paragraphs:
            text += para.text + "\n"
    elif suffix in ["png", "jpg", "jpeg"]:
        image = Image.open(uploaded_file)
        text = pytesseract.image_to_string(image)
    elif suffix == "csv":
        df = pd.read_csv(uploaded_file)
        text = df.to_string()
    elif suffix in ["xls", "xlsx"]:
        df = pd.read_excel(uploaded_file)
        text = df.to_string()
    else:
        text = uploaded_file.read().decode("utf-8", errors="ignore")

    return text.strip()

def create_embeddings(texts, model_name="all-MiniLM-L6-v2"):
    """Convert text chunks to embeddings."""
    model = SentenceTransformer(model_name)
    embeddings = model.encode(texts, convert_to_numpy=True)
    return embeddings

def chunk_text(text, chunk_size=500):
    """Split long text into smaller chunks."""
    words = text.split()
    return [" ".join(words[i:i + chunk_size]) for i in range(0, len(words), chunk_size)]


uploaded_files = st.file_uploader(
    "Upload PDF, PPTX, DOCX, PNG, JPG, CSV, XLSX files",
    type=["pdf", "pptx", "docx", "png", "jpg", "jpeg", "csv", "xls", "xlsx"],
    accept_multiple_files=True
)

if uploaded_files:
    all_texts = []
    file_refs = []

    for uploaded_file in uploaded_files:
        text = extract_text_from_file(uploaded_file)
        chunks = chunk_text(text)
        for chunk in chunks:
            all_texts.append(chunk)
            file_refs.append(uploaded_file.name)


    embeddings = create_embeddings(all_texts)
    dim = embeddings.shape[1]
    index = faiss.IndexFlatL2(dim)
    index.add(np.array(embeddings))

    st.success(f"Processed {len(uploaded_files)} files and indexed {len(all_texts)} text chunks.")


    query = st.text_input("Enter a question about your documents:")

    if query:
        query_vec = create_embeddings([query])
        D, I = index.search(query_vec, k=3)

        st.subheader("Retrieved Answers")

        for i, idx in enumerate(I[0]):
            result_text = all_texts[idx]
            source_file = file_refs[idx]

           
            highlighted_text = result_text
            for word in query.split():
                highlighted_text = highlighted_text.replace(
                    word, f"<span class='highlight'>{word}</span>"
                )

            st.markdown(f"**Result {i+1}:**", unsafe_allow_html=True)
            st.markdown(highlighted_text, unsafe_allow_html=True)
            st.markdown(f"**Source File:** `{source_file}`", unsafe_allow_html=True)

   
            col1, col2 = st.columns([1, 1])
            with col1:
                st.markdown(
                    f"<a class='source-link' href='#' target='_blank'>🔗 View Original</a>",
                    unsafe_allow_html=True
                )
            with col2:
            
                with open(source_file, "rb") as f:
                    st.download_button(
                        label="⬇️ Download File",
                        data=f,
                        file_name=Path(source_file).name,
                        key=f"download_{uuid.uuid4()}" 
                    )

else:
    st.info("Please upload one or more files to start.")

"""What does the document say about AI in climate change?”
“List the algorithms mentioned in the PowerPoint.”
“Who are the employees listed in the Excel sheet?”
“What does the DOCX say about generative AI?”
“What is written in the image?”
"""
